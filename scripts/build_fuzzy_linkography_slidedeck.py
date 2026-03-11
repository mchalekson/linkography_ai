#!/usr/bin/env python3
from __future__ import annotations

from pathlib import Path

import matplotlib
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

matplotlib.use("Agg")
import matplotlib.pyplot as plt

from linkography_ai.fuzzy_linkography import fuzzy_weight_matrix, load_chunk_moves, semantic_similarity
from linkography_ai.paths import data_v2_root


REPO_ROOT = Path(__file__).resolve().parents[1]
DOCS_DIR = REPO_ROOT / "docs"
TABLES_DIR = REPO_ROOT / "outputs" / "tables"
ANALYSIS_DIR = REPO_ROOT / "outputs" / "analysis"
TMP_FIG_DIR = DOCS_DIR / "_generated_fuzzy_figures"
OUT_PPTX = DOCS_DIR / "FUZZY_LINKOGRAPHY_V2_SLIDEDECK.pptx"


def _read_inputs() -> tuple[pd.DataFrame, pd.DataFrame, pd.DataFrame, pd.DataFrame, pd.DataFrame]:
    outcomes = pd.read_csv(TABLES_DIR / "fuzzy_linkography_outcomes_tests.csv")
    robustness = pd.read_csv(TABLES_DIR / "fuzzy_linkography_conference_robustness.csv")
    examples = pd.read_csv(TABLES_DIR / "fuzzy_linkography_example_meetings.csv")
    increment = pd.read_csv(TABLES_DIR / "fuzzy_linkography_model_increment.csv")
    merged_session = pd.read_csv(TABLES_DIR / "fuzzy_linkography_with_outcomes_by_session.csv")
    return outcomes, robustness, examples, increment, merged_session


def _save_feature_pvalues(outcomes: pd.DataFrame, out_path: Path) -> None:
    df = outcomes.copy()
    df["label"] = df["feature"].str.replace("_", "\n", regex=False)
    colors = ["#B23A48" if p < 0.05 else "#7A8793" for p in df["p_mannwhitney"]]
    plt.figure(figsize=(8.8, 5.1))
    plt.bar(df["label"], df["p_mannwhitney"], color=colors)
    plt.axhline(0.05, linestyle="--", color="#111827", linewidth=1)
    plt.ylabel("p-value")
    plt.title("Direct Outcome Tests for Derived Trajectory Measures")
    plt.tight_layout()
    plt.savefig(out_path, dpi=200)
    plt.close()


def _save_conference_diff(robustness: pd.DataFrame, out_path: Path) -> None:
    df = robustness[robustness["feature"] == "mean_nonzero_weight"].copy()
    plt.figure(figsize=(8.4, 4.8))
    colors = ["#2E8B57" if x >= 0 else "#B23A48" for x in df["mean_diff"]]
    plt.bar(df["conference"], df["mean_diff"], color=colors)
    plt.axhline(0.0, color="#111827", linewidth=1)
    plt.ylabel("Funded mean - unfunded mean")
    plt.title("Conference Robustness: Semantic-Link Strength Direction")
    plt.tight_layout()
    plt.savefig(out_path, dpi=200)
    plt.close()


def _save_model_increment(increment: pd.DataFrame, out_path: Path) -> None:
    df = increment.copy()
    label_map = {
        "baseline_full_plus_fuzzy": "Current funding model\n+ v2 meeting layer",
        "baseline_full": "Current funding model",
    }
    df = df[df["model"].isin(label_map.keys())].copy()
    df["label"] = df["model"].map(label_map)
    plt.figure(figsize=(8.4, 4.8))
    plt.bar(df["label"], df["cv_auc_mean"], color=["#4C78A8", "#B23A48"])
    plt.ylabel("CV AUC")
    plt.ylim(0.45, max(0.85, float(df["cv_auc_mean"].max()) + 0.05))
    plt.xticks(rotation=15, ha="right")
    plt.title("Funding-Model Check")
    plt.tight_layout()
    plt.savefig(out_path, dpi=200)
    plt.close()


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_bullets_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(22)


def _add_picture_slide(prs: Presentation, title: str, image_path: Path, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(str(image_path), Inches(0.45), Inches(1.1), width=Inches(6.2))
    textbox = slide.shapes.add_textbox(Inches(6.9), Inches(1.2), Inches(2.5), Inches(4.5))
    tf = textbox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)


def _add_table_slide(
    prs: Presentation,
    title: str,
    df: pd.DataFrame,
    *,
    left: float = 0.45,
    top: float = 1.2,
    width: float = 9.0,
    height: float = 3.4,
    font_size: int = 14,
    note: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    rows, cols = df.shape[0] + 1, df.shape[1]
    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table

    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True

    for i, row in enumerate(df.itertuples(index=False), start=1):
        for j, value in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(value)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)

    if note:
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top + height + 0.1), Inches(width), Inches(0.45))
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(10)


def _build_worked_example() -> tuple[pd.DataFrame, pd.DataFrame, str, str]:
    root = data_v2_root()
    meeting_dir = root / "2020NES" / "output_2020_11_05_NES_S1" / "1_DAC_Simulations_Zoom_Meeting_2020_11_05_12_24_10"
    chunk_ids = [1, 3, 5]
    move_rows = []
    moves = []
    for chunk_i in chunk_ids:
        chunk_fp = meeting_dir / f"1_DAC_Simulations_Zoom_Meeting_2020_11_05_12_24_10_chunk{chunk_i}.json"
        # Fallback to direct JSON load for nested structure
        import json
        raw = json.loads(chunk_fp.read_text())
        cs = raw.get("chunk_summary", {})
        chunk_moves = load_chunk_moves(chunk_fp, chunk_i)
        if not chunk_moves:
            continue
        move = chunk_moves[0]
        moves.append(move)
        excerpt = move.text.strip().replace("\n", " ")
        move_rows.append(
            {
                "Chunk": chunk_i,
                "Time": f"{int(move.start_sec//60):02d}:{int(move.start_sec%60):02d}",
                "idea_trajectory": cs.get("idea_trajectory", ""),
                "decision level": cs.get("decision_crystallization_level", ""),
                "full utterance quote": excerpt,
            }
        )

    link_rows = []
    if moves:
        texts = [m.text for m in moves]
        sim = semantic_similarity(texts, method="lsa")
        weights = fuzzy_weight_matrix(sim, threshold=0.35)
        for i in range(len(moves)):
            for j in range(i + 1, len(moves)):
                w = float(weights[i, j])
                if w > 0:
                    link_rows.append(
                        {
                            "From": f"Chunk {moves[i].chunk_index}",
                            "To": f"Chunk {moves[j].chunk_index}",
                            "Weight": f"{w:.3f}",
                        }
                    )
    if not link_rows:
        link_rows.append({"From": "No links", "To": "above threshold", "Weight": "0.000"})

    note = (
        "* Worked example from one canonical SCIALOG meeting. Decision level comes from chunk_summary "
        "(1 = earlier/open discussion, 3 = more crystallized discussion). The left table comes directly from the JSON content; "
        "the right table is added by fuzzy linkography, which infers weighted links between those utterances across the meeting. "
        "Each row is one utterance from a longer meeting dialogue, not the full dialogue."
    )
    header = "Conference: 2020NES | Session: 2020_11_05_NES_S1 | Meeting: 1_DAC_Simulations_Zoom_Meeting_2020_11_05_12_24_10"
    return pd.DataFrame(move_rows), pd.DataFrame(link_rows), note, header


def _add_worked_example_slide(prs: Presentation) -> None:
    json_df, link_df, note, header = _build_worked_example()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Worked Example: Added Meeting-Level Layer"
    header_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.75), Inches(9.0), Inches(0.35))
    header_tf = header_box.text_frame
    header_tf.word_wrap = True
    header_tf.paragraphs[0].text = header
    header_tf.paragraphs[0].font.size = Pt(11)
    left_rows, left_cols = json_df.shape[0] + 1, json_df.shape[1]
    left = slide.shapes.add_table(left_rows, left_cols, Inches(0.25), Inches(1.1), Inches(7.1), Inches(2.9)).table
    for j, col in enumerate(json_df.columns):
        left.cell(0, j).text = str(col)
        for p in left.cell(0, j).text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.bold = True
    for i, row in enumerate(json_df.itertuples(index=False), start=1):
        for j, value in enumerate(row):
            left.cell(i, j).text = str(value)
            for p in left.cell(i, j).text_frame.paragraphs:
                p.font.size = Pt(7)

    right_rows, right_cols = link_df.shape[0] + 1, link_df.shape[1]
    right = slide.shapes.add_table(right_rows, right_cols, Inches(7.45), Inches(1.1), Inches(2.0), Inches(1.45)).table
    for j, col in enumerate(link_df.columns):
        right.cell(0, j).text = str(col)
        for p in right.cell(0, j).text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.bold = True
    for i, row in enumerate(link_df.itertuples(index=False), start=1):
        for j, value in enumerate(row):
            right.cell(i, j).text = str(value)
            for p in right.cell(i, j).text_frame.paragraphs:
                p.font.size = Pt(8)

    note_box = slide.shapes.add_textbox(Inches(0.25), Inches(4.1), Inches(9.2), Inches(0.95))
    tf = note_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = note
    p.font.size = Pt(10)


def build_deck() -> Path:
    TMP_FIG_DIR.mkdir(parents=True, exist_ok=True)
    outcomes, robustness, examples, increment, merged_session = _read_inputs()

    pvals_fig = TMP_FIG_DIR / "feature_pvalues.png"
    robustness_fig = TMP_FIG_DIR / "conference_diff.png"
    increment_fig = TMP_FIG_DIR / "model_increment.png"
    _save_feature_pvalues(outcomes, pvals_fig)
    _save_conference_diff(robustness, robustness_fig)
    _save_model_increment(increment, increment_fig)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    top_model = increment.sort_values("cv_auc_mean", ascending=False).iloc[0]
    base_model = increment[increment["model"] == "baseline_full"].iloc[0]
    delta = float(top_model["cv_auc_mean"] - base_model["cv_auc_mean"])
    direct_top = outcomes.iloc[0]
    conf_focus = robustness[robustness["feature"] == "mean_nonzero_weight"].copy()
    n_consistent = int((conf_focus["mean_diff"] < 0).sum())
    n_conf_robustness = int(conf_focus["conference"].nunique())
    n_meetings = int(pd.read_csv(TABLES_DIR / "fuzzy_linkography_v2_by_meeting.csv").shape[0])
    n_chunks = int(pd.read_csv(TABLES_DIR / "fuzzy_linkography_v2_by_chunk.csv").shape[0])
    n_sessions_with_outcomes = int(merged_session["any_funded"].notna().sum())
    conference_list = sorted(merged_session["conference"].dropna().astype(str).unique().tolist())

    coverage_df = pd.DataFrame(
        [
            ["Conferences included", f"{len(conference_list)} ({', '.join(conference_list)})"],
            ["Meetings analyzed", str(n_meetings)],
            ["Chunks analyzed", str(n_chunks)],
            ["Sessions with outcomes", str(n_sessions_with_outcomes)],
            ["Methodology", "Fuzzy linkography"],
            ["Similarity model", "Latent Semantic Analysis"],
        ],
        columns=["Metric", "Value"],
    )
    source_vs_method_df = pd.DataFrame(
        [
            ["chunk_summary.idea_trajectory", "weighted links inferred between utterances across the full meeting"],
            ["chunk_summary.decision_crystallization_level", "meeting-level summary scores built from those inferred links"],
            ["chunk_summary.collective_engagement_level", "cross-chunk connection structure across the whole session"],
            ["utterance_annotations", "a session-level layer that is not directly stored in the JSON files"],
        ],
        columns=["Directly from Evey's JSON outputs", "Added by fuzzy linkography"],
    )
    v2_connection_df = pd.DataFrame(
        [
            ["idea_trajectory", "The chunk-level direction of discussion in the JSONs", "This added layer looks across the whole meeting and asks how strongly utterances stay related over time"],
            ["decision_crystallization_level", "How clearly the discussion is moving toward a decision", "This added layer tracks whether later discussion reconnects to earlier content as decisions develop"],
            ["collective_engagement_level", "How active and mutually engaged the group seems", "This added layer can show whether continuity is mostly within one speaker or carried across speakers"],
            ["shared_vision / commitment flags", "Signals that the team is aligning around an idea", "This added layer helps show whether that alignment is sustained across the meeting rather than only inside one chunk"],
        ],
        columns=["v2 JSON term", "What it already gives us", "What the added meeting-level layer contributes"],
    )

    conf_table = conf_focus.loc[:, ["conference", "n_sessions", "funded_mean", "unfunded_mean", "mean_diff"]].copy()
    conf_table["funded_mean"] = conf_table["funded_mean"].map(lambda x: f"{x:.3f}")
    conf_table["unfunded_mean"] = conf_table["unfunded_mean"].map(lambda x: f"{x:.3f}")
    conf_table["mean_diff"] = conf_table["mean_diff"].map(lambda x: f"{x:+.3f}")
    conf_table.columns = ["Conference", "N", "Funded Mean", "Unfunded Mean", "Diff"]

    example_cols = ["conference", "session_id_norm", "n_chunks", "any_funded"]
    high_examples = examples[examples["example_group"] == "high"].loc[:, example_cols].head(4).copy()
    low_examples = examples[examples["example_group"] == "low"].loc[:, example_cols].head(4).copy()
    for df in [high_examples, low_examples]:
        df["any_funded"] = df["any_funded"].map(lambda x: "NA" if pd.isna(x) else str(int(x)))
        df.columns = ["Conference", "Session", "Chunks", "Funded"]

    increment_table = increment[increment["model"].isin(["baseline_full", "baseline_full_plus_fuzzy"])].loc[:, ["model", "cv_auc_mean"]].copy()
    increment_table["model"] = increment_table["model"].map(
        {
            "baseline_full_plus_fuzzy": "Current funding model + v2 meeting layer",
            "baseline_full": "Current funding model",
        }
    ).fillna(increment_table["model"])
    increment_table["cv_auc_mean"] = increment_table["cv_auc_mean"].map(lambda x: f"{x:.3f}")
    increment_table.columns = ["Model", "CV AUC"]

    _add_title_slide(
        prs,
        "Semantic Trajectory Analysis on Canonical v2 Outputs",
        "Weighted semantic links across meetings, using the updated v2 annotation schema",
    )
    _add_bullets_slide(
        prs,
        "Why This Layer",
        [
            "The canonical v2 outputs already track idea trajectory, decision crystallization, and related meeting-state signals.",
            "This layer uses a fuzzy-linkography approach to add weighted links between utterances across the meeting.",
            "The current implementation uses Latent Semantic Analysis as the semantic-similarity model inside that approach.",
            "The current rerun also filters out auxiliary ATTN files.",
        ],
    )
    _add_table_slide(prs, "Coverage", coverage_df, width=8.3, height=3.2, left=0.6, top=1.0, font_size=15)
    _add_table_slide(
        prs,
        "What Comes From The JSONs vs What We Compute",
        source_vs_method_df,
        top=1.0,
        height=2.6,
        font_size=14,
        note="* The JSON files already contain chunk-level annotations. The fuzzy-linkography step adds a separate meeting-level layer by inferring weighted semantic links across utterances over time.",
    )
    _add_worked_example_slide(prs)
    _add_bullets_slide(
        prs,
        "Current Direct Comparison",
        [
            "When the goal is funding, this added meeting-level layer does not give one clear score that separates funded from unfunded sessions.",
            "Some differences appear in individual comparisons, but they do not stay consistent enough to treat as the main finding.",
            "So the main thing this analysis is pulling out right now is session-level structure, not a clean funding signal.",
        ],
    )
    _add_table_slide(prs, "How This Connects To The v2 JSON Terms", v2_connection_df, top=1.0, height=3.2, font_size=13)
    _add_picture_slide(
        prs,
        "Across Conferences",
        robustness_fig,
        [
            f"One direction repeats in {n_consistent} of {n_conf_robustness} conferences, but not across the full corpus.",
            "Direction is mixed rather than stable across cohorts.",
            "So this does not support one robust conference-invariant funding story.",
        ],
    )
    _add_table_slide(prs, "Conference Breakdown", conf_table, top=1.15, height=2.45, font_size=15)
    example_note = (
        "* Review examples only: these meetings were selected from the highest and lowest "
        "mean_nonzero_weight values after filtering to meetings with at least 10 extracted moves. "
        "They are not 'best' or 'worst' meetings."
    )
    _add_table_slide(
        prs,
        "Example Meetings To Review: Higher-Score Cases",
        high_examples,
        top=1.2,
        height=2.1,
        font_size=16,
        note=example_note,
    )
    _add_table_slide(
        prs,
        "Example Meetings To Review: Lower-Score Cases",
        low_examples,
        top=1.2,
        height=2.1,
        font_size=16,
        note=example_note,
    )
    _add_picture_slide(
        prs,
        "Main Modeling Result",
        increment_fig,
        [
            "This slide checks one question: does adding the v2 meeting-level layer help the current funding model?",
            f"Current funding model AUC = {base_model['cv_auc_mean']:.3f}.",
            f"Current funding model + v2 meeting layer = {increment[increment['model'] == 'baseline_full_plus_fuzzy']['cv_auc_mean'].iloc[0]:.3f}.",
            f"Estimated change = {increment[increment['model'] == 'baseline_full_plus_fuzzy']['cv_auc_mean'].iloc[0] - base_model['cv_auc_mean']:.3f}.",
        ],
    )
    _add_bullets_slide(
        prs,
        "What This Means Right Now",
        [
            "The canonical-source rerun gives broader coverage and cleaner outcome linkage.",
            "But in this version, the new layer from Evey's updated JSON outputs is not improving the existing outcome model.",
            "So the current value is to summarize and compare whole meetings on top of chunk_summary and utterance_annotations, rather than to claim better funding prediction yet.",
        ],
    )
    _add_table_slide(prs, "Simple Model Comparison", increment_table, top=1.45, height=1.6, font_size=18, width=7.5, left=1.0)
    _add_bullets_slide(
        prs,
        "How This Fits The v2 Outputs",
        [
            "Evey's canonical outputs add the updated JSON schema across 8 conferences, including chunk_summary and utterance_annotations.",
            "This added meeting-level layer is a way to summarize the full session on top of those updated chunk annotations.",
            "At the moment, its strongest use is to help characterize meetings and select sessions for closer review, not to claim a predictive gain by itself.",
        ],
    )
    _add_bullets_slide(
        prs,
        "Next Step",
        [
            "Read the example meetings back against their chunk_summary and utterance_annotations.",
            "Check whether the meeting-level patterns line up with idea_trajectory, decision_crystallization_level, and shared_vision_indicator in the underlying JSONs.",
            "If needed, revise the semantic-similarity component while keeping the same overall fuzzy-linkography framework.",
            "Keep the interpretation in the v2 annotation language rather than introducing a separate vocabulary.",
        ],
    )

    prs.save(OUT_PPTX)
    return OUT_PPTX


if __name__ == "__main__":
    out = build_deck()
    print(f"Saved: {out}")
